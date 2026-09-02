#!/usr/bin/env python3
"""Four distinct 8.5x11 portrait flyer options for FURM.
Output: furm_flyers.pptx  (4 slides = 4 options)
Run: .pptxvenv/bin/python make_flyers.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BARK  = RGBColor(0x3d,0x2b,0x1f)
BARK2 = RGBColor(0x5a,0x42,0x30)
CREAM = RGBColor(0xf6,0xef,0xe3)
CREAM2= RGBColor(0xfb,0xf6,0xec)
PINE  = RGBColor(0x2f,0x5d,0x43)
PINE2 = RGBColor(0x3f,0x7a,0x58)
MOSS  = RGBColor(0x6f,0x9c,0x6a)
SUN   = RGBColor(0xe6,0xb2,0x3e)
SUN2  = RGBColor(0xf2,0xc7,0x66)
GOLDT = RGBColor(0x7a,0x5b,0x16)
WHITE = RGBColor(0xff,0xff,0xff)
SOFT  = RGBColor(0xEC,0xE6,0xD6)

LOGO_CIRCLE = "../assets/logo_circle.png"
LOGO = "../assets/logo_padded.png"
QR = "../assets/qr_placeholder.png"
FONT = "Georgia"
W, H = 8.5, 11.0

def logo_src():
    return LOGO_CIRCLE if os.path.exists(LOGO_CIRCLE) else LOGO

def bg(s, color=CREAM):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

def rect(s, l, t, w, h, color, rounded=True, line=None, lw=1.5):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shp, Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh

def oval(s, l, t, d, fill, line=None, lw=3):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l),Inches(t),Inches(d),Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh

def badge(s, cx, cy, d, ring=SUN):
    oval(s, cx-d/2, cy-d/2, d, CREAM2, line=ring, lw=3)
    src = logo_src()
    if os.path.exists(src):
        pd = d*0.96
        s.shapes.add_picture(src, Inches(cx-pd/2), Inches(cy-pd/2), width=Inches(pd), height=Inches(pd))

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=None):
    tb = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,(txt,size,color,bold,italic) in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if sp is not None: p.space_after=Pt(sp)
        r=p.add_run(); r.text=txt
        f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name=FONT
    return tb

def rich(s, l, t, w, h, segs, size, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    p=tf.paragraphs[0]; p.alignment=align
    for txt,opt in segs:
        r=p.add_run(); r.text=txt
        f=r.font; f.size=Pt(size); f.name=FONT
        f.bold=opt.get("bold",False); f.italic=opt.get("italic",False)
        f.underline=opt.get("underline",False); f.color.rgb=opt.get("color",BARK2)
    return tb

def qr_block_dark(s, l, t, w, h, note1, note2):
    rect(s, l, t, w, h, PINE, rounded=True)
    if os.path.exists(QR):
        s.shapes.add_picture(QR, Inches(l+0.2), Inches(t+(h-1.15)/2), width=Inches(1.15), height=Inches(1.15))
    text(s, l+1.45, t, w-1.55, h, [
        ("Scan To Join", 18, WHITE, True, False),
        (note1, 12.5, SUN2, False, False),
        (note2, 12.5, SOFT, False, True),
    ], anchor=MSO_ANCHOR.MIDDLE, sp=2)

EVENTS = [
    ("Kickoff Social", "Late September", "Food, Games & Meeting Your People"),
    ("CubCare Mentor / Mentee", "October", "Get Paired With Someone Who Gets It"),
    ("Life @ Med School", "November", "An Honest Panel With Med Students"),
    ("De-Stress Events", "Reading Period", "Comfort Food & Calm During Finals"),
    ("Giveaways", "All Semester!", "Free Stuff, Thanks To Our Sponsors"),
]

def new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================ OPTION A
# Bold top-heavy: large pine header, giant headline, event list, QR footer
def option_a(prs):
    s=new(prs); bg(s)
    rect(s,-0.2,-0.2,W+0.4,3.0,PINE,rounded=False)
    rect(s,-0.2,2.75,W+0.4,0.12,SUN,rounded=False)
    badge(s, W/2, 1.25, 2.1, ring=SUN)
    text(s,0.4,2.28,W-0.8,0.5,[("FIRST-GENERATION UNDERREPRESENTED IN RESEARCH & MEDICINE",11,SUN2,True,False)],align=PP_ALIGN.CENTER)
    text(s,0.5,3.15,W-1.0,1.0,[("You Belong Here.",42,PINE,True,False)],align=PP_ALIGN.CENTER)
    rich(s,0.7,4.2,W-1.4,0.9,[
        ("However ",{"color":BARK2}),("YOU",{"bold":True,"underline":True,"color":PINE}),
        (" Define First-Gen Or Underrepresented \u2014 Come Find Your People.",{"color":BARK2}),
    ],15,align=PP_ALIGN.CENTER)
    text(s,0.9,5.15,W-1.8,0.4,[("This Fall At FURM",18,PINE,True,False)],align=PP_ALIGN.CENTER)
    y=5.7
    for name,when,blurb in EVENTS:
        rich(s,1.1,y,W-2.2,0.5,[
            ("\u2022  ",{"color":SUN,"bold":True}),
            (name+" \u00b7 "+when,{"color":BARK,"bold":True}),
            ("  \u2014  "+blurb,{"color":BARK2,"italic":True}),
        ],12.5,align=PP_ALIGN.LEFT)
        y+=0.44
    text(s,0.9,8.05,W-1.8,0.4,[("Interested In Leadership? Open E-Board Roles \u2014 Priority Deadline Sept 18.",11.5,PINE,True,True)],align=PP_ALIGN.CENTER)
    qr_block_dark(s,0.9,8.6,W-1.8,1.9,"Our Mailing List, Mentorship & E-Board","Brown Bears \u2014 Come As You Are.")

# ============================================================ OPTION B
# Editorial minimal: lots of whitespace, big centered logo, few words, QR
def option_b(prs):
    s=new(prs); bg(s, CREAM2)
    rect(s,0,0,W,0.35,PINE,rounded=False)
    rect(s,0,H-0.35,W,0.35,PINE,rounded=False)
    badge(s, W/2, 2.5, 2.7, ring=PINE)
    text(s,0.5,4.15,W-1.0,0.5,[("FURM \u00b7 Brown PLME",15,GOLDT,True,False)],align=PP_ALIGN.CENTER)
    text(s,0.5,4.7,W-1.0,1.2,[
        ("You Belong Here.",40,PINE,True,False),
    ],align=PP_ALIGN.CENTER)
    rich(s,1.0,5.85,W-2.0,1.0,[
        ("First-Gen? Underrepresented In Medicine Or Research? However ",{"color":BARK2}),
        ("YOU",{"bold":True,"underline":True,"color":PINE}),
        (" Define That \u2014 There's A Place For You Here.",{"color":BARK2}),
    ],16,align=PP_ALIGN.CENTER)
    # thin divider
    rect(s,3.25,7.05,2.0,0.04,SUN,rounded=False)
    text(s,0.8,7.25,W-1.6,0.9,[
        ("Community \u00b7 Mentorship \u00b7 Belonging",15,PINE,True,False),
        ("Events All Semester, Plus Our CubCare Mentorship Program.",13,BARK2,False,True),
    ],align=PP_ALIGN.CENTER,sp=4)
    # centered QR
    if os.path.exists(QR):
        s.shapes.add_picture(QR, Inches(W/2-0.85), Inches(8.35), width=Inches(1.7), height=Inches(1.7))
    text(s,0.5,10.05,W-1.0,0.5,[("Scan To Join Our Mailing List & Programs",14,PINE,True,False)],align=PP_ALIGN.CENTER)

# ============================================================ OPTION C
# Two-tone vertical split: pine left rail with logo + identity, content right
def option_c(prs):
    s=new(prs); bg(s)
    railw=3.1
    rect(s,-0.2,-0.2,railw+0.2,H+0.4,PINE,rounded=False)
    rect(s,railw,-0.2,0.12,H+0.4,SUN,rounded=False)
    badge(s, railw/2, 1.6, 2.2, ring=SUN)
    text(s,0.25,2.95,railw-0.5,3.0,[
        ("First-Generation",18,WHITE,True,False),
        ("Underrepresented",18,WHITE,True,False),
        ("In Research &",18,WHITE,True,False),
        ("Medicine",18,WHITE,True,False),
    ],align=PP_ALIGN.CENTER,sp=2)
    text(s,0.25,5.5,railw-0.5,2.0,[
        ("Community",15,SUN2,True,False),
        ("Care",15,SUN2,True,False),
        ("Belonging",15,SUN2,True,False),
    ],align=PP_ALIGN.CENTER,sp=6)
    text(s,0.25,9.6,railw-0.5,1.0,[("Brown Bears \u2014\nCome As You Are.",12,SOFT,False,True)],align=PP_ALIGN.CENTER)
    # right content
    cx=railw+0.5; cw=W-railw-0.9
    text(s,cx,0.7,cw,1.0,[("You Belong Here.",34,PINE,True,False)])
    rich(s,cx,1.75,cw,1.0,[
        ("However ",{"color":BARK2}),("YOU",{"bold":True,"underline":True,"color":PINE}),
        (" Define First-Gen Or Underrepresented \u2014 Come Find Your People.",{"color":BARK2}),
    ],15)
    text(s,cx,2.95,cw,0.4,[("This Fall",18,PINE,True,False)])
    y=3.45
    for name,when,blurb in EVENTS:
        text(s,cx,y,cw,0.6,[
            (name+" \u00b7 "+when,13,BARK,True,False),
            (blurb,11.5,BARK2,False,True),
        ],sp=1)
        y+=0.72
    qr_block_dark(s,cx,8.9,cw,1.6,"Mailing List, Mentorship & E-Board","Priority E-Board Deadline: Sept 18")

# ============================================================ OPTION D
# Poster: big headline, gold hook band, 2-col event cards, QR footer
def option_d(prs):
    s=new(prs); bg(s)
    badge(s, W/2, 1.15, 1.9, ring=PINE)
    text(s,0.5,2.15,W-1.0,0.9,[("You Belong Here.",40,PINE,True,False)],align=PP_ALIGN.CENTER)
    rich(s,0.7,3.05,W-1.4,0.7,[
        ("However ",{"color":BARK2}),("YOU",{"bold":True,"underline":True,"color":PINE}),
        (" Define It \u2014 Come Find Your People.",{"color":BARK2}),
    ],15,align=PP_ALIGN.CENTER)
    # gold hook band
    rect(s,0.8,3.75,W-1.6,0.85,SUN,rounded=True)
    text(s,0.9,3.8,W-1.8,0.75,[
        ("Find Your People.",21,BARK,True,False),
        ("A Home For First-Gen & Underrepresented Students In Medicine",11.5,GOLDT,False,True),
    ],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=1)
    # 2-col event cards (5 events -> 3 top, 2 bottom + leadership)
    cards = EVENTS + [("Join The E-Board","Sept 18","Open Leadership Roles This Fall")]
    cw=(W-1.6-0.2)/2; ch=1.05
    x0=0.8; y0=4.85
    for i,(name,when,blurb) in enumerate(cards):
        col=i%2; row=i//2
        x=x0+col*(cw+0.2); y=y0+row*(ch+0.15)
        rect(s,x,y,cw,ch,CREAM2,line=PINE,lw=1.4)
        text(s,x+0.12,y+0.08,cw-0.24,ch-0.16,[
            (name,13.5,PINE,True,False),
            (when,10.5,GOLDT,True,True),
            (blurb,10.5,BARK2,False,True),
        ],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=0)
    yqr=y0+3*(ch+0.15)+0.05
    qr_block_dark(s,0.8,yqr,W-1.6,1.55,"Mailing List, Mentorship & E-Board","Brown Bears \u2014 Come As You Are.")

def main():
    prs=Presentation()
    prs.slide_width=Inches(W); prs.slide_height=Inches(H)
    option_a(prs); option_b(prs); option_c(prs); option_d(prs)
    prs.save("furm_flyers.pptx")
    print("wrote furm_flyers.pptx (4 options: A top-heavy, B editorial, C split-rail, D poster)")

main()
