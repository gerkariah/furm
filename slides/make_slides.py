#!/usr/bin/env python3
"""Generates FURM PowerPoint decks, redesigned to match the forest/earthy
website + link tree aesthetic. Produces multiple design options.

  furm_slides.pptx  - 4 slides: 2 'about' concepts + 2 'flyer' concepts

Run: .pptxvenv/bin/python make_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- palette (matches site) ----
BARK  = RGBColor(0x3d,0x2b,0x1f)
BARK2 = RGBColor(0x5a,0x42,0x30)
CREAM = RGBColor(0xf6,0xef,0xe3)
CREAM2= RGBColor(0xfb,0xf6,0xec)
PINE  = RGBColor(0x2f,0x5d,0x43)
PINE2 = RGBColor(0x3f,0x7a,0x58)
MOSS  = RGBColor(0x6f,0x9c,0x6a)
SUN   = RGBColor(0xe6,0xb2,0x3e)
SUN2  = RGBColor(0xf2,0xc7,0x66)
GOLDT = RGBColor(0x7a,0x5b,0x16)
WHITE = RGBColor(0xff,0xff,0xff)
SOFT  = RGBColor(0xEC,0xE6,0xD6)

LOGO = "../assets/logo_padded.png"
LOGO_CIRCLE = "../assets/logo_circle.png"
QR = "../assets/qr_placeholder.png"
FONT_HEAD = "Georgia"
FONT_BODY = "Georgia"

EW, EH = Inches(13.333), Inches(7.5)

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i,(txt,size,color,bold,italic) in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        if sp is not None: p.space_after = Pt(sp)
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = FONT_HEAD
    return tb

def rich(slide, l, t, w, h, segments, size, base_color, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """One paragraph made of multiple runs.
    segments = list of (text, {bold,italic,underline,color,caps})"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for txt, opt in segments:
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(size); f.name = FONT_HEAD
        f.bold = opt.get("bold", False)
        f.italic = opt.get("italic", False)
        f.underline = opt.get("underline", False)
        f.color.rgb = opt.get("color", base_color)
    return tb

def rect(slide, l, t, w, h, color, rounded=True, line=None, lw=1.5):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def oval(slide, l, t, d, fill, line=None, lw=3):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l),Inches(t),Inches(d),Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def logo_badge(slide, cx, cy, d, ring=PINE):
    """circular ring + the pre-cropped circular logo, both truly round"""
    oval(slide, cx-d/2, cy-d/2, d, CREAM2, line=ring, lw=3)
    src = LOGO_CIRCLE if os.path.exists(LOGO_CIRCLE) else LOGO
    if os.path.exists(src):
        pd = d*0.96
        slide.shapes.add_picture(src, Inches(cx-pd/2), Inches(cy-pd/2),
                                 width=Inches(pd), height=Inches(pd))

def bg(slide, color=CREAM):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def sub_pill(slide, l, t, txt, w=2.4):
    rect(slide, l, t, w, 0.42, PINE, rounded=True)
    text(slide, l, t-0.02, w, 0.46, [(txt, 13, WHITE, True, False)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# --------------------------------------------------------------------------
# SLIDE A1 — "About" concept 1: logo left, warm and clean
# --------------------------------------------------------------------------
def about_1(prs):
    s = blank(prs); bg(s)
    rect(s, -0.2, -0.2, 13.733, 7.9, CREAM, rounded=False)
    # left pine panel + sun divider
    rect(s, -0.3, -0.3, 4.7, 8.1, PINE, rounded=False)
    rect(s, 4.35, -0.3, 0.14, 8.1, SUN, rounded=False)
    logo_badge(s, 2.05, 2.1, 3.0, ring=SUN)
    text(s, 0.3, 3.85, 3.6, 0.6, [("Brown PLME", 15, SUN2, True, False)], align=PP_ALIGN.CENTER)
    text(s, 0.3, 4.2, 3.6, 1.6, [
        ("First-Generation", 20, WHITE, True, False),
        ("Underrepresented in", 20, WHITE, True, False),
        ("Research & Medicine", 20, WHITE, True, False),
    ], align=PP_ALIGN.CENTER, sp=2)
    text(s, 0.3, 6.35, 3.6, 0.6, [("Community \u00b7 Care \u00b7 Belonging", 13, SOFT, False, True)], align=PP_ALIGN.CENTER)

    # right content
    text(s, 5.0, 0.55, 7.9, 0.9, [("Who we are", 34, PINE, True, False)])
    # paragraph 1 with YOU emphasis (caps, bold, underline, italic)
    you = {"bold":True,"italic":True,"underline":True,"color":PINE}
    base = {"color":BARK2}
    rich(s, 5.0, 1.45, 7.9, 1.5, [
        ("A home within Brown PLME for students who care about community, mentorship, and equity "
         "in healthcare. We're first-generation and underrepresented in research and medicine \u2014 and ", base),
        ("YOU", you),
        (" get to define what that means for ", base),
        ("YOU", you),
        (".", base),
    ], 16, BARK2)
    text(s, 5.0, 2.95, 7.9, 0.9, [
        ("Whoever you are, if you believe in building more welcoming healthcare spaces, there's a "
         "place for you here \u2014 we're about belonging and community.", 16, BARK2, False, False),
    ])

    # fall events 2x2 grid + Giveaways row, QR to the right
    text(s, 5.0, 3.95, 5.0, 0.45, [("This fall", 22, PINE, True, False)])
    ev = [("Kickoff Social","Coming\nlate September"),("CubCare Mentorship","October"),
          ("Life @ Med School","November"),("De-Stress Night","Reading Period")]
    gx, gy, cw, ch = 5.0, 4.5, 2.55, 0.85
    for i,(name,when) in enumerate(ev):
        col = i % 2; row = i // 2
        x = gx + col*(cw+0.18); y = gy + row*(ch+0.16)
        rect(s, x, y, cw, ch, CREAM2, line=PINE, lw=1.4)
        text(s, x+0.06, y+0.05, cw-0.12, ch-0.1, [
            (name, 13, PINE, True, False),
            (when.replace("\n"," "), 10.5, GOLDT, False, True),
        ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp=1)
    # giveaways full-width chip under the grid
    gyv = gy + 2*(ch+0.16)
    rect(s, gx, gyv, cw*2+0.18, 0.72, SUN, rounded=True)
    text(s, gx+0.1, gyv+0.03, cw*2-0.02, 0.66, [
        ("Giveaways \u00b7 All semester!", 13.5, BARK, True, False),
        ("Sponsors: AWS, Elks Foundation, Millennium Fellowship", 10, GOLDT, False, True),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp=0)

    # QR block on the right
    qx, qy = 10.5, 4.5
    rect(s, qx, qy, 2.35, 2.55, PINE, rounded=True)
    if os.path.exists(QR):
        s.shapes.add_picture(QR, Inches(qx+0.55), Inches(qy+0.25), width=Inches(1.25), height=Inches(1.25))
    text(s, qx+0.1, qy+1.55, 2.15, 1.0, [
        ("Scan to join", 15, WHITE, True, False),
        ("mailing list \u00b7 mentorship \u00b7 e-board", 11, SUN2, False, True),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, sp=2)

# --------------------------------------------------------------------------
# SLIDE A2 — "About" concept 2: centered, badge on top, event cards
# --------------------------------------------------------------------------
def about_2(prs):
    s = blank(prs); bg(s)
    rect(s, -0.2, -0.2, 13.733, 2.0, PINE, rounded=False)
    rect(s, -0.2, 1.75, 13.733, 0.12, SUN, rounded=False)
    logo_badge(s, 1.5, 1.0, 2.1, ring=SUN)
    text(s, 2.9, 0.32, 9.8, 1.4, [
        ("FURM", 40, WHITE, True, False),
        ("First-Generation Underrepresented in Research & Medicine  \u00b7  Brown PLME", 15, SUN2, False, True),
    ], anchor=MSO_ANCHOR.MIDDLE, sp=2)
    text(s, 0.8, 2.15, 11.7, 1.7, [
        ("We're a home within Brown PLME for students who care about community, mentorship, "
         "and equity in healthcare. We're first-generation and underrepresented in research and "
         "medicine \u2014 and YOU define what that means for you.", 17, BARK2, False, False),
        ("Whoever you are, there's a place for you here. We're about belonging and community.",
         17, PINE, True, True),
    ], align=PP_ALIGN.CENTER, sp=8)
    # four event cards
    cards = [("Kickoff Social","Coming late September","Food, games & meeting your people"),
             ("CubCare Mentorship","October","Get paired with someone who gets it"),
             ("Life @ Med School","November","An honest panel with med students"),
             ("De-Stress Night","Reading period","Comfort food during finals")]
    x = 0.7; w = 2.95
    for name,when,desc in cards:
        rect(s, x, 4.15, w, 2.15, CREAM2, line=PINE, lw=1.6)
        rect(s, x+0.35, 4.02, w-0.7, 0.4, SUN, rounded=True)
        text(s, x+0.35, 4.0, w-0.7, 0.44, [(when, 11, BARK, True, False)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x+0.12, 4.6, w-0.24, 1.6, [
            (name, 17, PINE, True, False),
            (desc, 12.5, BARK2, False, False),
        ], align=PP_ALIGN.CENTER, sp=6)
        x += w + 0.13
    text(s, 0.7, 6.55, 11.9, 0.7, [
        ("Scan the QR to join our mailing list & sign up for mentorship  \u00b7  we do giveaways too!",
         14, PINE, True, False)], align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------------
# SLIDE F1 — "Flyer" concept 1: big warm invite + QR
# --------------------------------------------------------------------------
def flyer_1(prs):
    s = blank(prs); bg(s)
    logo_badge(s, 6.666, 1.7, 2.6, ring=PINE)
    text(s, 1.0, 3.15, 11.33, 0.7, [("You belong here.", 40, PINE, True, False)],
         align=PP_ALIGN.CENTER)
    text(s, 1.5, 3.95, 10.33, 1.0, [
        ("First-generation? Underrepresented in medicine or research? "
         "However YOU define that \u2014 come find your people.", 20, BARK2, False, False),
    ], align=PP_ALIGN.CENTER)
    # kickoff banner
    rect(s, 2.2, 5.0, 8.93, 1.2, SUN, rounded=True)
    text(s, 2.3, 5.12, 8.73, 1.0, [
        ("KICKOFF SOCIAL  \u00b7  Coming late September", 24, BARK, True, False),
        ("Food \u00b7 games \u00b7 community \u2014 watch your email for the date", 15, GOLDT, False, True),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp=2)
    text(s, 1.0, 6.5, 11.33, 0.6, [
        ("Scan to join our mailing list & mentorship   \u00b7   Brown Bears \u2014 come as you are",
         16, PINE, True, False)], align=PP_ALIGN.CENTER)
    # QR on the right
    rect(s, 11.0, 4.95, 1.55, 1.55, CREAM2, line=PINE, lw=1.5)
    if os.path.exists(QR):
        s.shapes.add_picture(QR, Inches(11.15), Inches(5.1), width=Inches(1.25), height=Inches(1.25))

# --------------------------------------------------------------------------
# SLIDE F2 — "Flyer" concept 2: two-column, events + leadership
# --------------------------------------------------------------------------
def flyer_2(prs):
    s = blank(prs); bg(s)
    rect(s, -0.2, -0.2, 13.733, 1.6, PINE, rounded=False)
    rect(s, -0.2, 1.35, 13.733, 0.12, SUN, rounded=False)
    logo_badge(s, 1.35, 0.72, 1.7, ring=SUN)
    text(s, 2.6, 0.1, 10.4, 1.25, [
        ("You belong here.", 30, WHITE, True, False),
        ("FURM \u00b7 First-Gen & Underrepresented in Research & Medicine \u00b7 Brown PLME", 14, SUN2, False, True),
    ], anchor=MSO_ANCHOR.MIDDLE, sp=2)
    text(s, 0.8, 1.65, 11.7, 0.9, [
        ("However YOU define first-generation or underrepresented \u2014 come find your people. "
         "We're about community, care, and belonging.", 17, BARK2, False, False),
    ], align=PP_ALIGN.CENTER)
    # left: fall events
    rect(s, 0.7, 2.75, 5.9, 3.55, CREAM2, line=PINE, lw=1.6)
    text(s, 0.95, 2.9, 5.4, 0.5, [("Fall Events", 22, PINE, True, False)])
    events = [("Kickoff Social","Coming late September"),
              ("CubCare Mentor / Mentee Program","October"),
              ("Life @ Med School (honest panel)","November"),
              ("De-Stress Night","Reading Period"),
              ("Giveaways","All semester!")]
    y = 3.5
    for name,when in events:
        text(s, 1.05, y, 5.4, 0.5, [
            ("\u2022  "+name, 16, BARK, True, False),
            ("      "+when, 12.5, GOLDT, False, True),
        ], sp=1)
        y += 0.56
    # right: leadership
    rect(s, 6.75, 2.75, 5.85, 3.55, PINE, rounded=True)
    text(s, 7.0, 2.9, 5.4, 0.5, [("Interested in leadership?", 20, WHITE, True, False)])
    text(s, 7.0, 3.4, 5.4, 0.5, [("Open E-Board positions this fall (year-long):", 13.5, SUN2, False, True)])
    roles = ["Secretary","Treasurer","Underclassman Liaison (\u201929 & \u201930)","Upperclassman Liaison (\u201927 & \u201928)"]
    y = 3.95
    for rname in roles:
        text(s, 7.1, y, 5.3, 0.45, [("\u2022  "+rname, 15.5, WHITE, False, False)])
        y += 0.5
    text(s, 7.0, 6.0, 5.4, 0.3, [("Priority deadline: September 18", 13.5, SUN2, True, True)])
    # footer link + qr
    text(s, 0.7, 6.6, 9.4, 0.6, [
        ("Scan to join our mailing list, mentorship & E-Board \u2014 all in one place", 15, PINE, True, False)],
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 11.0, 6.4, 1.0, 0.95, CREAM2, line=PINE, lw=1.5)
    if os.path.exists(QR):
        s.shapes.add_picture(QR, Inches(11.08), Inches(6.47), width=Inches(0.85), height=Inches(0.85))

# --------------------------------------------------------------------------
# PORTRAIT PAPER FLYER — 8.5 x 11 in
# --------------------------------------------------------------------------
def portrait_flyer(prs):
    s = blank(prs); bg(s)
    W = 8.5
    # top pine header band
    rect(s, -0.2, -0.2, W+0.4, 2.55, PINE, rounded=False)
    rect(s, -0.2, 2.3, W+0.4, 0.12, SUN, rounded=False)
    logo_badge(s, W/2, 1.15, 2.0, ring=SUN)

    # headline
    text(s, 0.5, 2.75, W-1.0, 0.9, [("You Belong Here.", 40, PINE, True, False)], align=PP_ALIGN.CENTER)
    you = {"bold":True,"italic":False,"underline":True,"color":PINE}
    base = {"color":BARK2}
    rich(s, 0.6, 3.7, W-1.2, 1.0, [
        ("First-Generation? Underrepresented In Medicine Or Research? However ", base),
        ("YOU", you),
        (" Define That \u2014 Come Find Your People.", base),
    ], 16, BARK2, align=PP_ALIGN.CENTER)

    # gold banner — the emotional hook (not an event; events are listed below)
    rect(s, 0.9, 4.7, W-1.8, 1.0, SUN, rounded=True)
    text(s, 1.0, 4.77, W-2.0, 0.86, [
        ("Find Your People.", 24, BARK, True, False),
        ("A Home For First-Gen & Underrepresented Students In Medicine", 12.5, GOLDT, False, True),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp=2)

    # what's this fall (name + short blurb per event, Title Case)
    text(s, 0.9, 5.95, W-1.8, 0.45, [("This Fall At FURM", 19, PINE, True, False)], align=PP_ALIGN.CENTER)
    items = [
        ("Kickoff Social \u00b7 Late September", "Food, Games & Meeting Your People"),
        ("CubCare Mentor / Mentee Program \u00b7 October", "Get Paired With Someone Who Gets It"),
        ("Life @ Med School \u00b7 November", "An Honest Panel With Med Students"),
        ("De-Stress Events \u00b7 Reading Period", "Comfort Food & Calm During Finals"),
        ("Giveaways \u00b7 All Semester!", "Free Stuff, Thanks To Our Sponsors"),
    ]
    y = 6.45
    for name, blurb in items:
        rich(s, 1.2, y, W-2.4, 0.5, [
            ("\u2022  ", {"color":SUN,"bold":True}),
            (name, {"color":BARK,"bold":True}),
            ("   \u2014  "+blurb, {"color":BARK2,"italic":True}),
        ], 12.5, BARK2)
        y += 0.42

    # E-board recruiting line
    text(s, 0.9, 8.55, W-1.8, 0.4, [
        ("Interested In Leadership? We Have Open E-Board Positions \u2014 Priority Deadline Sept 18.", 12, PINE, True, True)
    ], align=PP_ALIGN.CENTER)

    # QR + link footer
    rect(s, 0.9, 9.05, W-1.8, 1.55, PINE, rounded=True)
    if os.path.exists(QR):
        s.shapes.add_picture(QR, Inches(1.2), Inches(9.25), width=Inches(1.15), height=Inches(1.15))
    text(s, 2.6, 9.2, W-3.1, 1.3, [
        ("Scan To Join", 19, WHITE, True, False),
        ("Our Mailing List, Mentorship & E-Board", 12.5, SUN2, False, False),
        ("Brown Bears \u2014 Come As You Are.", 12.5, SOFT, False, True),
    ], anchor=MSO_ANCHOR.MIDDLE, sp=2)

def main():
    prs = Presentation()
    prs.slide_width = EW; prs.slide_height = EH
    about_1(prs)
    about_2(prs)
    flyer_1(prs)
    flyer_2(prs)
    prs.save("furm_slides.pptx")
    print("wrote furm_slides.pptx (4 slides: about-1, about-2, flyer-1, flyer-2)")

    # NOTE: the portrait flyer is NOT regenerated here.
    # The final flyer is the hand-edited Option D in furm_flyers.pptx — do not overwrite it.

main()
